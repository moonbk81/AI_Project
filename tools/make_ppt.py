import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip('#')
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

prs = Presentation()

# 테마 색상 설정 (다크 모드 & 라임 포인트)
BG_COLOR = hex_to_rgb('000000')
ACCENT_COLOR = hex_to_rgb('deff9a')
TEXT_COLOR = hex_to_rgb('f8fafc')
SUB_TEXT_COLOR = hex_to_rgb('94a3b8')

def apply_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # 하단 라임 포인트 라인
    line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(7.1), Inches(9), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.visible = False

# ==========================================
# 1. Title Slide
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide)

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Intelligent Telephony"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = TEXT_COLOR
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Log Analysis System"
p2.font.size = Pt(54)
p2.font.bold = True
p2.font.color.rgb = ACCENT_COLOR
p2.alignment = PP_ALIGN.CENTER

p3 = tf.add_paragraph()
p3.text = "\nProgress Report: Achieving 90.9% Accuracy via Local RAG Optimization"
p3.font.size = Pt(20)
p3.font.color.rgb = SUB_TEXT_COLOR
p3.alignment = PP_ALIGN.CENTER

# ==========================================
# 2. Benchmark Slide
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide)

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title.text_frame.text = "Local LLM Benchmarking (VRAM 8GB)"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

rows, cols = 4, 3
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(3)).table

headers = ["Model Name", "Latency", "Routing Accuracy"]
for c, h in enumerate(headers):
    cell = table.cell(0, c)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = hex_to_rgb('1a1a1a')
    cell.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

data = [
    ["Gemma4:e4b (Q4) [최종 선정]", "40.56s", "100% (완벽)"],
    ["Qwen2.5:7b", "19.98s", "86.1% (누락 발생)"],
    ["Gemma3:4b", "113.02s", "83.3% (OOM 발생)"]
]

for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        cell = table.cell(r+1, c)
        cell.text = val
        cell.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR

# ==========================================
# 3. Optimization Leap (Before & After)
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide)

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title.text_frame.text = "Optimization Performance Leap"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

# Before 박스
box1 = slide.shapes.add_shape(1, Inches(1), Inches(2.5), Inches(3.5), Inches(3))
box1.fill.solid()
box1.fill.fore_color.rgb = hex_to_rgb('1a1a1a')
box1.line.color.rgb = hex_to_rgb('ef4444') # 빨간색

b_text = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(3.5), Inches(2))
b_p = b_text.text_frame.paragraphs[0]
b_p.text = "Initial Score\n17.2%"
b_p.font.size = Pt(40)
b_p.font.bold = True
b_p.font.color.rgb = hex_to_rgb('ef4444')
b_p.alignment = PP_ALIGN.CENTER

# After 박스 (강조)
box2 = slide.shapes.add_shape(1, Inches(5.5), Inches(2.5), Inches(3.5), Inches(3))
box2.fill.solid()
box2.fill.fore_color.rgb = hex_to_rgb('111827')
box2.line.color.rgb = ACCENT_COLOR

a_text = slide.shapes.add_textbox(Inches(5.5), Inches(2.7), Inches(3.5), Inches(2))
a_p = a_text.text_frame.paragraphs[0]
a_p.text = "Optimized Score\n90.9%"
a_p.font.size = Pt(44)
a_p.font.bold = True
a_p.font.color.rgb = ACCENT_COLOR
a_p.alignment = PP_ALIGN.CENTER

# ==========================================
# 4. Governance
# ==========================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_background(slide)

title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title.text_frame.text = "Security & Quality Governance"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR

for i, (title_txt, desc_txt) in enumerate([
    ("100% On-Premise", "사내망 로그 데이터 유출\n원천 차단 (Zero-Leakage)"),
    ("LLM-as-a-Judge", "Gemma3:12b 모델을 통한\n환각(Hallucination) 통제")]):

    tile = slide.shapes.add_shape(1, Inches(1 + i*4.5), Inches(2.5), Inches(3.5), Inches(2.5))
    tile.fill.solid()
    tile.fill.fore_color.rgb = hex_to_rgb('1a1a1a')
    tile.line.color.rgb = hex_to_rgb('334155')

    t_box = slide.shapes.add_textbox(Inches(1 + i*4.5), Inches(3), Inches(3.5), Inches(1))
    tf = t_box.text_frame
    p1 = tf.paragraphs[0]
    p1.text = title_txt
    p1.font.size = Pt(24)
    p1.font.color.rgb = ACCENT_COLOR
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = desc_txt
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_COLOR
    p2.alignment = PP_ALIGN.CENTER

# 바탕화면에 저장
desktop_path = os.path.join(os.path.expanduser('~'), 'Desktop', 'Log_Analysis_Report_Styled.pptx')
prs.save(desktop_path)
print(f"PPT 파일 생성 완료: {desktop_path}")